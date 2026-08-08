"""
UniformDesk — System Analyst presentation builder.
Generates a comprehensive .pptx with live screenshots, flows, and simple animations.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

ROOT = Path(__file__).resolve().parent
SHOTS = ROOT / "screenshots"
OUT = ROOT / "UniformDesk_System_Overview.pptx"

# Brand palette
NAVY = RGBColor(0x0B, 0x1F, 0x1C)
TEAL = RGBColor(0x0F, 0x6E, 0x56)
TEAL_SOFT = RGBColor(0xE6, 0xF4, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x5C, 0x6B, 0x66)
ACCENT = RGBColor(0x1A, 0x8F, 0x6E)
WARN = RGBColor(0xB5, 0x6A, 0x00)
DARK_CARD = RGBColor(0x12, 0x2A, 0x26)


def set_run(run, size=18, bold=False, color=NAVY, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill=TEAL, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    # softer corners via adj if available
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def add_text_box(slide, l, t, w, h, text, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_bullets(slide, l, t, w, h, items, size=16, color=NAVY):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        set_run(run, size=size, color=color)
    return box


def enable_transition(slide, duration_ms=400):
    """Add a simple fade transition via OOXML."""
    sld = slide._element
    # remove existing transition
    for child in list(sld):
        if child.tag == qn("p:transition"):
            sld.remove(child)
    tr = etree.SubElement(sld, qn("p:transition"))
    tr.set("spd", "med")
    fade = etree.SubElement(tr, qn("p:fade"))
    return tr


def add_appear_animation(slide, shape, order=1):
    """Simple fade-in entrance animation for a shape (click-advance)."""
    sld = slide._element
    # Ensure timing/tnLst structure exists under cSld sibling <p:timing>
    timing = sld.find(qn("p:timing"))
    if timing is None:
        timing = etree.SubElement(sld, qn("p:timing"))
    tn_lst = timing.find(qn("p:tnLst"))
    if tn_lst is None:
        tn_lst = etree.SubElement(timing, qn("p:tnLst"))
    # One root par sequence
    root_par = tn_lst.find(qn("p:par"))
    if root_par is None:
        root_par = etree.SubElement(tn_lst, qn("p:par"))
        ctn = etree.SubElement(root_par, qn("p:cTn"))
        ctn.set("id", "1")
        ctn.set("dur", "indefinite")
        ctn.set("restart", "never")
        ctn.set("nodeType", "tmRoot")
        child_tn = etree.SubElement(ctn, qn("p:childTnLst"))
        seq = etree.SubElement(child_tn, qn("p:seq"))
        seq.set("concurrent", "1")
        seq.set("nextAc", "seek")
        seq_ctn = etree.SubElement(seq, qn("p:cTn"))
        seq_ctn.set("id", "2")
        seq_ctn.set("dur", "indefinite")
        seq_ctn.set("nodeType", "mainSeq")
        etree.SubElement(seq_ctn, qn("p:childTnLst"))
        prev = etree.SubElement(seq, qn("p:prevCondLst"))
        prev_cond = etree.SubElement(prev, qn("p:cond"))
        prev_cond.set("evt", "onPrev")
        prev_tgt = etree.SubElement(prev_cond, qn("p:tgtEl"))
        etree.SubElement(prev_tgt, qn("p:sldTgt"))
        nxt = etree.SubElement(seq, qn("p:nextCondLst"))
        nxt_cond = etree.SubElement(nxt, qn("p:cond"))
        nxt_cond.set("evt", "onNext")
        nxt_tgt = etree.SubElement(nxt_cond, qn("p:tgtEl"))
        etree.SubElement(nxt_tgt, qn("p:sldTgt"))

    # Navigate to mainSeq childTnLst
    root_par = tn_lst.find(qn("p:par"))
    seq = root_par.find(".//" + qn("p:seq"))
    main_child = seq.find(qn("p:cTn")).find(qn("p:childTnLst"))

    anim_id = str(100 + order)
    par = etree.SubElement(main_child, qn("p:par"))
    ctn = etree.SubElement(par, qn("p:cTn"))
    ctn.set("id", anim_id)
    ctn.set("fill", "hold")
    st = etree.SubElement(ctn, qn("p:stCondLst"))
    cond = etree.SubElement(st, qn("p:cond"))
    cond.set("delay", "0" if order == 1 else "250")
    child = etree.SubElement(ctn, qn("p:childTnLst"))
    anim_par = etree.SubElement(child, qn("p:par"))
    anim_ctn = etree.SubElement(anim_par, qn("p:cTn"))
    anim_ctn.set("id", str(200 + order))
    anim_ctn.set("presetID", "10")  # fade
    anim_ctn.set("presetClass", "entr")
    anim_ctn.set("presetSubtype", "0")
    anim_ctn.set("fill", "hold")
    anim_ctn.set("grpId", "0")
    anim_ctn.set("nodeType", "clickEffect")
    st2 = etree.SubElement(anim_ctn, qn("p:stCondLst"))
    cond2 = etree.SubElement(st2, qn("p:cond"))
    cond2.set("delay", "0")
    child2 = etree.SubElement(anim_ctn, qn("p:childTnLst"))
    anim_effect = etree.SubElement(child2, qn("p:animEffect"))
    anim_effect.set("transition", "in")
    anim_effect.set("filter", "fade")
    c_bhvr = etree.SubElement(anim_effect, qn("p:cBhvr"))
    c_tn = etree.SubElement(c_bhvr, qn("p:cTn"))
    c_tn.set("id", str(300 + order))
    c_tn.set("dur", "500")
    tgt = etree.SubElement(c_bhvr, qn("p:tgtEl"))
    sp = etree.SubElement(tgt, qn("p:spTgt"))
    # shape id from cNvPr
    sp_id = shape._element.nvSpPr.cNvPr.get("id")
    sp.set("spid", str(sp_id))


def add_footer(slide, page, total, title="UniformDesk System Overview"):
    add_text_box(
        slide,
        Inches(0.5),
        Inches(7.15),
        Inches(8),
        Inches(0.3),
        f"{title}  |  Confidential",
        size=10,
        color=MUTED,
    )
    add_text_box(
        slide,
        Inches(11.2),
        Inches(7.15),
        Inches(1.5),
        Inches(0.3),
        f"{page} / {total}",
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def section_banner(slide, kicker, title, subtitle=None):
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.35), fill=TEAL)
    add_text_box(slide, Inches(0.55), Inches(0.22), Inches(12), Inches(0.3), kicker, size=12, bold=True, color=TEAL_SOFT)
    add_text_box(slide, Inches(0.55), Inches(0.48), Inches(12), Inches(0.55), title, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.55), Inches(1.55), Inches(12), Inches(0.4), subtitle, size=14, color=MUTED)


def shot_or_placeholder(slide, name, l, t, w, h):
    path = SHOTS / f"{name}.png"
    if path.exists():
        slide.shapes.add_picture(str(path), l, t, width=w, height=h)
        return True
    ph = add_rect(slide, l, t, w, h, fill=TEAL_SOFT, line=TEAL)
    tf = ph.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = f"[Screenshot pending: {name}]"
    set_run(run, size=12, color=MUTED)
    return False


def screenshot_slide(prs, page, total, kicker, title, shot_name, caption, points=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    enable_transition(slide)
    section_banner(slide, kicker, title, caption)
    # image area
    shot_or_placeholder(slide, shot_name, Inches(0.45), Inches(2.05), Inches(8.4), Inches(4.85))
    # callouts
    card = add_rect(slide, Inches(9.1), Inches(2.05), Inches(3.7), Inches(4.85), fill=TEAL_SOFT)
    add_text_box(slide, Inches(9.35), Inches(2.25), Inches(3.3), Inches(0.4), "Analyst notes", size=14, bold=True, color=TEAL)
    if points:
        add_bullets(slide, Inches(9.35), Inches(2.75), Inches(3.3), Inches(3.8), points, size=13, color=NAVY)
    add_footer(slide, page, total)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Estimate total slides for footer (approx)
    # We'll assign page numbers sequentially
    pages = []

    def track(slide):
        pages.append(slide)
        return slide

    # ========== 1 Title ==========
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    add_bg(s, TEAL)
    add_text_box(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.4), "SYSTEM ANALYSIS  ·  PRODUCT WALKTHROUGH", size=14, bold=True, color=TEAL_SOFT)
    add_text_box(s, Inches(0.8), Inches(2.3), Inches(11.5), Inches(1.0), "UniformDesk", size=54, bold=True, color=WHITE)
    add_text_box(
        s,
        Inches(0.8),
        Inches(3.4),
        Inches(11),
        Inches(0.8),
        "Supplier supply · School issue · Proof for every student",
        size=22,
        color=TEAL_SOFT,
    )
    add_text_box(
        s,
        Inches(0.8),
        Inches(4.5),
        Inches(11),
        Inches(1.0),
        "Comprehensive system presentation with live application screenshots\nand end-to-end operational flows.",
        size=16,
        color=WHITE,
    )
    add_text_box(s, Inches(0.8), Inches(6.4), Inches(11), Inches(0.4), "Prepared as a System Analyst briefing  |  Local demo environment", size=12, color=TEAL_SOFT)

    # ========== 2 Agenda ==========
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    section_banner(s, "01  AGENDA", "What this briefing covers")
    cols = [
        ("Context", ["Business problem", "Stakeholders & roles", "Solution positioning"]),
        ("System design", ["Architecture", "Tenancy & RBAC", "End-to-end spine flow"]),
        ("Live product tour", ["School desk screens", "Supplier portal screens", "Trust & offline"]),
        ("Readiness", ["What is solid today", "Pre-production gaps", "Recommended next steps"]),
    ]
    x = 0.5
    for title, items in cols:
        add_rect(s, Inches(x), Inches(1.9), Inches(2.9), Inches(4.6), fill=TEAL_SOFT)
        add_text_box(s, Inches(x + 0.2), Inches(2.15), Inches(2.5), Inches(0.4), title, size=16, bold=True, color=TEAL)
        add_bullets(s, Inches(x + 0.2), Inches(2.7), Inches(2.5), Inches(3.5), items, size=13)
        x += 3.15

    # ========== 3 Problem ==========
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    section_banner(s, "02  PROBLEM", "Why schools and suppliers need UniformDesk")
    problems = [
        ("No proof of issue", "Uniforms leave the store without a signed, shareable record."),
        ("Stock without a ledger", "Counts drift; shortages are informal or hidden."),
        ("Broken supply chain", "Orders, deliveries, and invoices live in chat/spreadsheets."),
        ("Weak audit trail", "Auditors cannot correlate student, slip, DN, and payment."),
    ]
    y = 1.9
    for title, body in problems:
        add_rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(1.1), fill=TEAL_SOFT)
        add_text_box(s, Inches(0.75), Inches(y + 0.18), Inches(11.8), Inches(0.35), title, size=16, bold=True, color=TEAL)
        add_text_box(s, Inches(0.75), Inches(y + 0.55), Inches(11.8), Inches(0.4), body, size=14, color=NAVY)
        y += 1.25

    # ========== 4 Solution ==========
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    section_banner(s, "03  SOLUTION", "One modular monolith, two portals, one truth")
    add_rect(s, Inches(0.5), Inches(1.9), Inches(6.0), Inches(4.6), fill=TEAL)
    add_text_box(s, Inches(0.8), Inches(2.2), Inches(5.4), Inches(0.4), "School Desk", size=22, bold=True, color=WHITE)
    add_bullets(
        s,
        Inches(0.8),
        Inches(2.8),
        Inches(5.4),
        Inches(3.4),
        [
            "Issue with mandatory signature",
            "Stock balances + ledger",
            "Receive against supplier DN",
            "Student history & reports",
            "Public proof link for guardians",
        ],
        size=15,
        color=WHITE,
    )
    add_rect(s, Inches(6.8), Inches(1.9), Inches(6.0), Inches(4.6), fill=DARK_CARD)
    add_text_box(s, Inches(7.1), Inches(2.2), Inches(5.4), Inches(0.4), "Supplier Portal", size=22, bold=True, color=WHITE)
    add_bullets(
        s,
        Inches(7.1),
        Inches(2.8),
        Inches(5.4),
        Inches(3.4),
        [
            "Multi-school portfolio",
            "Orders → pack → dispatch",
            "Invoices + payment confirmation",
            "White-label branding",
            "Activity & notifications",
        ],
        size=15,
        color=WHITE,
    )

    # ========== 5 Stakeholders ==========
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    section_banner(s, "04  STAKEHOLDERS", "Roles and responsibilities")
    roles = [
        ("Storekeeper", "Issue, receive, stock ops"),
        ("School Admin", "Catalog, kits, users, integrations"),
        ("Auditor", "Reports & audit export"),
        ("Supplier Admin", "Catalog, branding, portfolio"),
        ("Supplier Staff", "Pack, dispatch, collect"),
        ("Guardian", "View public proof (no login)"),
    ]
    positions = [(0.5, 1.9), (4.6, 1.9), (8.7, 1.9), (0.5, 4.3), (4.6, 4.3), (8.7, 4.3)]
    for (title, body), (x, y) in zip(roles, positions):
        add_rect(s, Inches(x), Inches(y), Inches(3.8), Inches(2.0), fill=TEAL_SOFT)
        add_text_box(s, Inches(x + 0.25), Inches(y + 0.4), Inches(3.3), Inches(0.4), title, size=16, bold=True, color=TEAL)
        add_text_box(s, Inches(x + 0.25), Inches(y + 0.95), Inches(3.3), Inches(0.7), body, size=13, color=NAVY)

    # ========== 6 Architecture ==========
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    section_banner(s, "05  ARCHITECTURE", "Modular monolith — Next.js + Prisma")
    layers = [
        ("Presentation", "School desk UI  ·  Supplier portal  ·  Public proof  ·  PWA shell"),
        ("Application", "Server Actions  ·  REST (/api/v1)  ·  Edge middleware (JWT)"),
        ("Domain modules", "issue · inventory · supply · payments · reports · integrations"),
        ("Data", "Prisma ORM  ·  SQLite (local/CI)  ·  Postgres path for production"),
    ]
    y = 1.9
    for title, body in layers:
        add_rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(1.05), fill=TEAL_SOFT)
        add_text_box(s, Inches(0.75), Inches(y + 0.15), Inches(3.2), Inches(0.35), title, size=15, bold=True, color=TEAL)
        add_text_box(s, Inches(4.0), Inches(y + 0.3), Inches(8.4), Inches(0.5), body, size=14, color=NAVY)
        y += 1.2

    # ========== 7 E2E Flow ==========
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    section_banner(s, "06  SYSTEM FLOW", "End-to-end operational spine")
    steps = [
        ("1", "Catalog\nalign SKUs"),
        ("2", "Order\n(PO)"),
        ("3", "Pack &\nDispatch DN"),
        ("4", "School\nReceive"),
        ("5", "Issue +\nSignature"),
        ("6", "Proof +\nInvoice/Pay"),
    ]
    x = 0.45
    anim_order = 1
    for num, label in steps:
        circ = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.65), Inches(2.2), Inches(0.55), Inches(0.55))
        circ.fill.solid()
        circ.fill.fore_color.rgb = TEAL
        circ.line.fill.background()
        tf = circ.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = num
        set_run(run, size=14, bold=True, color=WHITE)
        card = add_rect(s, Inches(x), Inches(3.0), Inches(1.9), Inches(1.6), fill=TEAL_SOFT)
        add_text_box(s, Inches(x + 0.1), Inches(3.3), Inches(1.7), Inches(1.2), label, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_appear_animation(s, circ, order=anim_order)
        add_appear_animation(s, card, order=anim_order + 1)
        anim_order += 2
        if num != "6":
            add_text_box(s, Inches(x + 1.85), Inches(3.5), Inches(0.3), Inches(0.4), "→", size=20, bold=True, color=TEAL)
        x += 2.1
    add_bullets(
        s,
        Inches(0.5),
        Inches(5.0),
        Inches(12.3),
        Inches(1.8),
        [
            "SKU matching is mandatory: supplier product SKU → school catalog SKU.",
            "Stock only moves with ledger entries (receive / issue / void / adjust).",
            "No issue without signature; guardians verify via public /v/[token].",
            "Payments are settlement confirmation (cash/bank/sandbox), not a PSP product.",
        ],
        size=14,
    )

    # ========== 8 Flow swimlanes ==========
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    section_banner(s, "07  FLOW DETAIL", "Who does what across the loop")
    add_rect(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(1.5), fill=TEAL_SOFT)
    add_text_box(s, Inches(0.7), Inches(2.05), Inches(2.2), Inches(0.35), "SUPPLIER", size=14, bold=True, color=TEAL)
    add_text_box(s, Inches(3.0), Inches(2.05), Inches(9.5), Inches(1.1), "Create/fulfill order → Pack DN → Dispatch → Invoice → Confirm payment", size=15, color=NAVY)
    add_rect(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(1.5), fill=RGBColor(0xE8, 0xF0, 0xFF))
    add_text_box(s, Inches(0.7), Inches(3.75), Inches(2.2), Inches(0.35), "SCHOOL", size=14, bold=True, color=RGBColor(0x1D, 0x4E, 0x89))
    add_text_box(s, Inches(3.0), Inches(3.75), Inches(9.5), Inches(1.1), "Receive DN into stock → Issue kit with signature → Share proof → Review invoices", size=15, color=NAVY)
    add_rect(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.3), fill=RGBColor(0xFF, 0xF4, 0xE0))
    add_text_box(s, Inches(0.7), Inches(5.45), Inches(2.2), Inches(0.35), "GUARDIAN", size=14, bold=True, color=WARN)
    add_text_box(s, Inches(3.0), Inches(5.45), Inches(9.5), Inches(0.9), "Open /v/[token] — verify student, items, date — no UniformDesk account required", size=15, color=NAVY)

    # Live screenshot slides — school
    school_shots = [
        ("08  LIVE UI", "Login", "01-login", "Entry point for school and supplier tenants.", ["JWT session cookie", "Demo seed accounts", "Edge middleware gate"]),
        ("09  SCHOOL DESK", "Home / Desk overview", "02-school-home", "Daily operational command center.", ["Today stats", "Needs attention", "Issue CTA"]),
        ("10  SCHOOL DESK", "Issue desk", "03-school-issue", "Signature-gated kit issue workflow.", ["Student search", "Kit / lines", "Offline queue ready"]),
        ("11  SCHOOL DESK", "Stock balances", "04-school-stock", "On-hand by item and size.", ["Ledger-backed", "Adjust with reason", "Reorder CTA"]),
        ("12  SCHOOL DESK", "Activity timeline", "05-school-activity", "Audit trail with correlation IDs.", ["Issue / void / receive", "Support-friendly refs"]),
        ("13  SCHOOL DESK", "Notifications", "06-school-notifications", "Actionable attention feed.", ["Low stock", "Unpaid invoices", "Inbound DNs"]),
        ("14  SCHOOL DESK", "Deliveries", "07-school-deliveries", "Receive against supplier DN.", ["SKU match check", "Posts stock + ledger"]),
        ("15  SCHOOL DESK", "Orders", "08-school-orders", "School purchase orders to supplier.", ["Status tracking", "Links to DNs"]),
        ("16  SCHOOL DESK", "Invoices", "09-school-invoices", "Supplier bills visible to school.", ["Settlement status", "Printable sheet"]),
        ("17  SCHOOL DESK", "Low-stock reorder", "10-school-reorder", "Suggest PO lines from thresholds.", ["SKU → supplier product", "Unmatched listed"]),
        ("18  SCHOOL DESK", "Students roster", "11-school-students", "Admission-based student master.", ["CSV import", "History per student"]),
        ("19  SCHOOL DESK", "Desk search", "12-school-search", "Cross-entity findability.", ["Students, slips, supply docs"]),
        ("20  SCHOOL DESK", "Reports", "13-school-reports", "Issued today & shortages.", ["Audit CSV export", "Print report"]),
        ("21  SCHOOL ADMIN", "Catalog", "14-school-catalog", "Items and sizes (admin).", ["SKU is receive key"]),
        ("22  SCHOOL ADMIN", "Kits", "15-school-kits", "Issue bundles for speed.", ["Default qtys per item"]),
        ("23  SCHOOL ADMIN", "Users & roles", "16-school-users", "Desk accounts and RBAC.", ["Admin / storekeeper / auditor"]),
        ("24  SCHOOL ADMIN", "Integrations", "17-school-integrations", "School Master API + SSO.", ["Roster sync", "API keys"]),
        ("25  SCHOOL DESK", "Student history", "18-school-student-history", "Per-student issue trail.", ["Slips, shortages, voids"]),
    ]

    for kicker, title, shot, caption, notes in school_shots:
        track(screenshot_slide(prs, 0, 0, kicker, title, shot, caption, notes))

    # Slip + proof if present
    if (SHOTS / "19-school-slip.png").exists():
        track(
            screenshot_slide(
                prs,
                0,
                0,
                "26  TRUST",
                "Issue slip (audit record)",
                "19-school-slip",
                "Internal issue record; payment method + reference at desk.",
                ["Slip number", "No parent signature required", "Browser printable"],
            )
        )
    if (SHOTS / "20-public-proof.png").exists():
        track(
            screenshot_slide(
                prs,
                0,
                0,
                "27  TRUST",
                "Public guardian proof",
                "20-public-proof",
                "No login — tokenized verification page.",
                ["/v/[token]", "QR from slip", "Parent-friendly"],
            )
        )
    else:
        track(
            screenshot_slide(
                prs,
                0,
                0,
                "27  RESILIENCE",
                "Offline fallback",
                "20-offline",
                "PWA offline page when navigation fails.",
                ["Cached issue desk path", "Queue sync on reconnect"],
            )
        )

    supplier_shots = [
        ("28  SUPPLIER", "Supply home", "21-supplier-home", "Multi-school portfolio overview.", ["Open deliveries", "Unpaid invoices", "Attention strip"]),
        ("29  SUPPLIER", "Orders", "22-supplier-orders", "Inbound school POs.", ["Confirm / fulfill path"]),
        ("30  SUPPLIER", "Deliveries", "23-supplier-deliveries", "Pack and dispatch queue.", ["DN numbers", "Status pipeline"]),
        ("31  SUPPLIER", "Invoices", "24-supplier-invoices", "Billing and collection.", ["Payment confirmation"]),
        ("32  SUPPLIER", "Catalog", "25-supplier-catalog", "Products & SKUs.", ["Must match school SKUs"]),
        ("33  SUPPLIER", "Schools portfolio", "26-supplier-schools", "Linked campuses.", ["Filter by school"]),
        ("34  SUPPLIER", "Activity", "27-supplier-activity", "Supply timeline.", ["DN / INV / PAY corr IDs"]),
        ("35  SUPPLIER", "Notifications", "28-supplier-notifications", "Dispatch & collect prompts.", ["Packed DNs", "Open invoices"]),
        ("36  SUPPLIER", "Search", "29-supplier-search", "Find schools, SKUs, docs.", ["Top-bar search"]),
        ("37  SUPPLIER", "Branding", "30-supplier-branding", "White-label portal look.", ["Mark, color, contacts"]),
    ]
    for kicker, title, shot, caption, notes in supplier_shots:
        track(screenshot_slide(prs, 0, 0, kicker, title, shot, caption, notes))

    if (SHOTS / "31-supplier-delivery-detail.png").exists():
        track(
            screenshot_slide(
                prs,
                0,
                0,
                "38  SUPPLIER",
                "Delivery note detail (printable)",
                "31-supplier-delivery-detail",
                "Operational DN with print pack.",
                ["Lines + sign-off", "Dispatch / invoice actions"],
            )
        )
    if (SHOTS / "32-supplier-invoice-detail.png").exists():
        track(
            screenshot_slide(
                prs,
                0,
                0,
                "39  SUPPLIER",
                "Invoice detail + payment confirmation",
                "32-supplier-invoice-detail",
                "Record cash/bank or sandbox STK.",
                ["Not a payment rails product", "Marks paid when covered"],
            )
        )
    if (SHOTS / "33-supplier-order-detail.png").exists():
        track(
            screenshot_slide(
                prs,
                0,
                0,
                "39b  SUPPLIER",
                "Order detail",
                "33-supplier-order-detail",
                "Confirmed school PO with line items.",
                ["Create DN from order", "Cancel when needed", "School linkage"],
            )
        )

    # Data / trust
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    section_banner(s, "40  DATA & TRUST", "What makes UniformDesk auditable")
    add_bullets(
        s,
        Inches(0.6),
        Inches(2.0),
        Inches(12),
        Inches(4.8),
        [
            "StockBalance + StockLedgerEntry — every movement is explainable.",
            "IssueSlip requires signature image; void requires reason and restores stock.",
            "publicToken enables guardian verification without exposing the desk.",
            "Correlation IDs (slip / DN / INV / PAY) appear in activity feeds for support.",
            "Tenant isolation: schoolId / supplierId enforced in session and queries.",
            "RBAC separates storekeeper write paths from auditor read paths.",
        ],
        size=16,
    )

    # Offline
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    section_banner(s, "41  OFFLINE / PWA", "Issue continuity on flaky school networks")
    add_bullets(
        s,
        Inches(0.6),
        Inches(2.0),
        Inches(12),
        Inches(4.8),
        [
            "Service worker caches app shell (production).",
            "Desk home / issue page warm an IndexedDB roster snapshot.",
            "/issue-offline runs from cache when the server is unreachable.",
            "Queued issues sync via POST /api/v1/issue when back online.",
            "Local on-hand adjusts while offline so successive issues stay realistic.",
            "Scope is intentional: issue-only offline — receive/admin still need network.",
        ],
        size=16,
    )

    # Maturity
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    section_banner(s, "42  READINESS", "Analyst assessment")
    add_rect(s, Inches(0.5), Inches(1.9), Inches(6.0), Inches(4.6), fill=TEAL_SOFT)
    add_text_box(s, Inches(0.75), Inches(2.15), Inches(5.5), Inches(0.4), "Strong today", size=18, bold=True, color=TEAL)
    add_bullets(
        s,
        Inches(0.75),
        Inches(2.7),
        Inches(5.5),
        Inches(3.5),
        [
            "Clear domain model & dual portals",
            "Ledger + signature invariants (tested)",
            "Supply loop to payment confirmation",
            "CI on GitHub + middleware gate",
            "Pilot-ready UX (search, notices, print)",
        ],
        size=14,
    )
    add_rect(s, Inches(6.8), Inches(1.9), Inches(6.0), Inches(4.6), fill=RGBColor(0xFF, 0xF0, 0xE8))
    add_text_box(s, Inches(7.05), Inches(2.15), Inches(5.5), Inches(0.4), "Before live production", size=18, bold=True, color=WARN)
    add_bullets(
        s,
        Inches(7.05),
        Inches(2.7),
        Inches(5.5),
        Inches(3.5),
        [
            "Postgres (leave SQLite)",
            "Staging + backup/restore drill",
            "Secrets rotation / no demo passwords",
            "Observability (logs / uptime)",
            "Limited pilot: 1 school + 1 supplier",
        ],
        size=14,
    )

    # Demo script
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    section_banner(s, "43  DEMO SCRIPT", "10-minute live walkthrough")
    add_bullets(
        s,
        Inches(0.6),
        Inches(2.0),
        Inches(12),
        Inches(4.8),
        [
            "1) Login as store@greenfield.school / desk1234 — show desk home.",
            "2) Issue a kit with signature — open slip, show QR/share.",
            "3) Open public proof in a private window.",
            "4) Login as supply@uniformdesk.co — show deliveries & dispatch.",
            "5) School receive DN — stock rises; show activity correlation.",
            "6) Supplier invoice + mark paid (confirmation).",
            "7) Close with notifications / reorder / print packs if time.",
        ],
        size=16,
    )

    # Closing
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    add_bg(s, TEAL)
    add_text_box(s, Inches(0.8), Inches(2.3), Inches(11.5), Inches(0.9), "UniformDesk", size=44, bold=True, color=WHITE)
    add_text_box(
        s,
        Inches(0.8),
        Inches(3.3),
        Inches(11.5),
        Inches(1.2),
        "One spine: supply → stock → signed issue → proof → settlement confirmation.",
        size=20,
        color=TEAL_SOFT,
    )
    add_text_box(
        s,
        Inches(0.8),
        Inches(5.2),
        Inches(11.5),
        Inches(0.8),
        "Questions & discussion",
        size=18,
        bold=True,
        color=WHITE,
    )

    # Fix footers with real page numbers
    total = len(pages)
    for i, slide in enumerate(pages, start=1):
        # title slide already branded; still add discreet page on non-title if needed
        if i == 1 or i == total:
            continue
        add_footer(slide, i, total)

    prs.save(str(OUT))
    print(f"Wrote {OUT} ({total} slides)")


if __name__ == "__main__":
    build()
