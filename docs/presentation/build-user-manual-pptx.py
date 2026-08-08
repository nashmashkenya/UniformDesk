"""
UniformDesk — User Manual (live screenshots).
Dynamic PPT with fade transitions and click-to-reveal callouts.
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
SHOTS = ROOT / "screenshots"
OUT = ROOT / "UniformDesk_User_Manual.pptx"

NAVY = RGBColor(0x0B, 0x1F, 0x1C)
TEAL = RGBColor(0x0F, 0x6E, 0x56)
TEAL_SOFT = RGBColor(0xE6, 0xF4, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x5C, 0x6B, 0x66)
CREAM = RGBColor(0xF7, 0xFA, 0xF8)


def set_run(run, size=18, bold=False, color=NAVY, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill=TEAL, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def add_text_box(slide, l, t, w, h, text, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_bullets(slide, l, t, w, h, items, size=14, color=NAVY):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"•  {item}"
        set_run(run, size=size, color=color)
    return box


def enable_transition(slide):
    sld = slide._element
    for child in list(sld):
        if child.tag == qn("p:transition"):
            sld.remove(child)
    tr = etree.SubElement(sld, qn("p:transition"))
    tr.set("spd", "med")
    etree.SubElement(tr, qn("p:fade"))


def add_appear_animation(slide, shape, order=1):
    sld = slide._element
    timing = sld.find(qn("p:timing"))
    if timing is None:
        timing = etree.SubElement(sld, qn("p:timing"))
    tn_lst = timing.find(qn("p:tnLst"))
    if tn_lst is None:
        tn_lst = etree.SubElement(timing, qn("p:tnLst"))
    root_par = tn_lst.find(qn("p:par"))
    if root_par is None:
        root_par = etree.SubElement(tn_lst, qn("p:par"))
        ctn = etree.SubElement(root_par, qn("p:cTn"))
        ctn.set("id", "1")
        ctn.set("dur", "indefinite")
        ctn.set("restart", "never")
        ctn.set("nodeType", "tmRoot")
        child_tn = etree.SubElement(ctn, qn("p:childTnLst"))
        seq = etree.SubElement(child_tn, qn("p:seq"))
        seq.set("concurrent", "1")
        seq.set("nextAc", "seek")
        seq_ctn = etree.SubElement(seq, qn("p:cTn"))
        seq_ctn.set("id", "2")
        seq_ctn.set("dur", "indefinite")
        seq_ctn.set("nodeType", "mainSeq")
        etree.SubElement(seq_ctn, qn("p:childTnLst"))
        prev = etree.SubElement(seq, qn("p:prevCondLst"))
        prev_cond = etree.SubElement(prev, qn("p:cond"))
        prev_cond.set("evt", "onPrev")
        prev_tgt = etree.SubElement(prev_cond, qn("p:tgtEl"))
        etree.SubElement(prev_tgt, qn("p:sldTgt"))
        nxt = etree.SubElement(seq, qn("p:nextCondLst"))
        nxt_cond = etree.SubElement(nxt, qn("p:cond"))
        nxt_cond.set("evt", "onNext")
        nxt_tgt = etree.SubElement(nxt_cond, qn("p:tgtEl"))
        etree.SubElement(nxt_tgt, qn("p:sldTgt"))

    root_par = tn_lst.find(qn("p:par"))
    seq = root_par.find(".//" + qn("p:seq"))
    main_child = seq.find(qn("p:cTn")).find(qn("p:childTnLst"))
    par = etree.SubElement(main_child, qn("p:par"))
    ctn = etree.SubElement(par, qn("p:cTn"))
    ctn.set("id", str(100 + order))
    ctn.set("fill", "hold")
    st = etree.SubElement(ctn, qn("p:stCondLst"))
    cond = etree.SubElement(st, qn("p:cond"))
    cond.set("delay", "0" if order == 1 else "200")
    child = etree.SubElement(ctn, qn("p:childTnLst"))
    anim_par = etree.SubElement(child, qn("p:par"))
    anim_ctn = etree.SubElement(anim_par, qn("p:cTn"))
    anim_ctn.set("id", str(200 + order))
    anim_ctn.set("presetID", "10")
    anim_ctn.set("presetClass", "entr")
    anim_ctn.set("presetSubtype", "0")
    anim_ctn.set("fill", "hold")
    anim_ctn.set("grpId", "0")
    anim_ctn.set("nodeType", "clickEffect")
    st2 = etree.SubElement(anim_ctn, qn("p:stCondLst"))
    cond2 = etree.SubElement(st2, qn("p:cond"))
    cond2.set("delay", "0")
    child2 = etree.SubElement(anim_ctn, qn("p:childTnLst"))
    anim_effect = etree.SubElement(child2, qn("p:animEffect"))
    anim_effect.set("transition", "in")
    anim_effect.set("filter", "fade")
    c_bhvr = etree.SubElement(anim_effect, qn("p:cBhvr"))
    c_tn = etree.SubElement(c_bhvr, qn("p:cTn"))
    c_tn.set("id", str(300 + order))
    c_tn.set("dur", "450")
    tgt = etree.SubElement(c_bhvr, qn("p:tgtEl"))
    sp = etree.SubElement(tgt, qn("p:spTgt"))
    sp.set("spid", str(shape._element.nvSpPr.cNvPr.get("id")))


def add_footer(slide, page, total):
    add_text_box(
        slide,
        Inches(0.5),
        Inches(7.15),
        Inches(9),
        Inches(0.3),
        "UniformDesk User Manual  |  Live product screens",
        size=10,
        color=MUTED,
    )
    add_text_box(
        slide,
        Inches(11.2),
        Inches(7.15),
        Inches(1.5),
        Inches(0.3),
        f"{page} / {total}",
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def section_banner(slide, kicker, title, subtitle=None):
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.28), fill=TEAL)
    add_text_box(
        slide, Inches(0.55), Inches(0.18), Inches(12), Inches(0.28), kicker, size=11, bold=True, color=TEAL_SOFT
    )
    add_text_box(
        slide, Inches(0.55), Inches(0.42), Inches(12), Inches(0.5), title, size=26, bold=True, color=WHITE
    )
    if subtitle:
        add_text_box(
            slide, Inches(0.55), Inches(1.42), Inches(12), Inches(0.35), subtitle, size=13, color=MUTED
        )


def shot_or_placeholder(slide, name, l, t, w, h):
    path = SHOTS / f"{name}.png"
    if path.exists():
        slide.shapes.add_picture(str(path), l, t, width=w, height=h)
        return True
    ph = add_rect(slide, l, t, w, h, fill=TEAL_SOFT, line=TEAL)
    tf = ph.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = f"[Capture: {name}]"
    set_run(run, size=12, color=MUTED)
    return False


def screenshot_slide(prs, kicker, title, shot_name, caption, points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    enable_transition(slide)
    section_banner(slide, kicker, title, caption)
    shot_or_placeholder(slide, shot_name, Inches(0.4), Inches(1.95), Inches(8.5), Inches(4.95))
    card = add_rect(slide, Inches(9.15), Inches(1.95), Inches(3.7), Inches(4.95), fill=TEAL_SOFT)
    add_text_box(
        slide, Inches(9.4), Inches(2.15), Inches(3.25), Inches(0.35), "How to use", size=14, bold=True, color=TEAL
    )
    bullets = add_bullets(slide, Inches(9.4), Inches(2.6), Inches(3.25), Inches(4.0), points, size=12)
    add_appear_animation(slide, card, 1)
    add_appear_animation(slide, bullets, 2)
    return slide


def divider_slide(prs, chapter, title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    enable_transition(slide)
    add_bg(slide, TEAL)
    add_text_box(
        slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.4), chapter, size=14, bold=True, color=TEAL_SOFT
    )
    add_text_box(
        slide, Inches(0.8), Inches(2.7), Inches(11.5), Inches(0.8), title, size=36, bold=True, color=WHITE
    )
    y = 3.8
    for i, line in enumerate(lines):
        box = add_text_box(slide, Inches(0.8), Inches(y), Inches(11), Inches(0.4), line, size=16, color=TEAL_SOFT)
        add_appear_animation(slide, box, i + 1)
        y += 0.45
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slides: list = []

    def track(slide, footer=True):
        slides.append((slide, footer))
        return slide

    # Title
    s = track(prs.slides.add_slide(prs.slide_layouts[6]), footer=False)
    enable_transition(s)
    add_bg(s, TEAL)
    add_text_box(
        s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.4),
        "USER MANUAL  ·  LIVE PRODUCT GUIDE", size=14, bold=True, color=TEAL_SOFT,
    )
    add_text_box(s, Inches(0.8), Inches(2.15), Inches(11.5), Inches(0.9), "UniformDesk", size=52, bold=True, color=WHITE)
    add_text_box(
        s, Inches(0.8), Inches(3.2), Inches(11), Inches(0.7),
        "Supplier-led uniforms for Kenyan senior schools — stock, co-issue, reports & print",
        size=20, color=TEAL_SOFT,
    )
    add_text_box(
        s, Inches(0.8), Inches(4.3), Inches(11), Inches(1.0),
        "Comprehensive walkthrough with live screenshots from the running application.\n"
        "Roles · Create school · Co-issue · Still owed · Stock take · Printable reports",
        size=15, color=WHITE,
    )
    add_text_box(
        s, Inches(0.8), Inches(6.3), Inches(11), Inches(0.35),
        "Password for all demo accounts: desk1234", size=13, bold=True, color=TEAL_SOFT,
    )

    # How to use
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    section_banner(s, "01  HOW TO USE THIS MANUAL", "Dynamic slides · click to reveal tips")
    cards = [
        ("Live screens", "Every product screen was captured from the running UniformDesk app."),
        ("Click animations", "Callout cards fade in on click — present step by step."),
        ("Two portals", "Supplier owns the system; school reporters issue & stock take."),
        ("Print packs", "Reports, stock, DN, and invoices print from the browser."),
    ]
    for i, (t, body) in enumerate(cards):
        x = 0.5 + (i % 2) * 6.35
        y = 1.85 + (i // 2) * 2.4
        card = add_rect(s, Inches(x), Inches(y), Inches(6.0), Inches(2.15), fill=TEAL_SOFT)
        add_text_box(s, Inches(x + 0.3), Inches(y + 0.35), Inches(5.4), Inches(0.4), t, size=18, bold=True, color=TEAL)
        add_text_box(s, Inches(x + 0.3), Inches(y + 0.9), Inches(5.4), Inches(0.9), body, size=14, color=NAVY)
        add_appear_animation(s, card, i + 1)

    # Roles
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    section_banner(s, "02  ROLES & LOGINS", "Who signs in where")
    roles = [
        ("Supplier admin", "supply@uniformdesk.co", "Create schools, catalog, DN, invoices, co-issue, branding"),
        ("Supplier staff", "staff@uniformdesk.co", "Pack/dispatch, co-issue, reports (no branding)"),
        ("School reporter (GFS)", "report@greenfield.school", "Issue, stock take, receive DN, reports"),
        ("School reporter (RVA)", "report@riverside.school", "Same school desk for Riverside Academy"),
    ]
    for i, (role, email, desc) in enumerate(roles):
        y = 1.7 + i * 1.2
        card = add_rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(1.05), fill=CREAM if i % 2 == 0 else TEAL_SOFT)
        add_text_box(s, Inches(0.75), Inches(y + 0.15), Inches(3.8), Inches(0.35), role, size=16, bold=True, color=TEAL)
        add_text_box(s, Inches(4.6), Inches(y + 0.15), Inches(4.5), Inches(0.35), email, size=14, bold=True, color=NAVY)
        add_text_box(s, Inches(0.75), Inches(y + 0.55), Inches(11.5), Inches(0.35), desc, size=13, color=MUTED)
        add_appear_animation(s, card, i + 1)

    # Flow
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    section_banner(s, "03  END-TO-END FLOW", "From create school to printed reports")
    steps = [
        ("1", "Create / link school", "Supplier admin"),
        ("2", "Catalog SKUs", "Match school items"),
        ("3", "Order → DN → receive", "Stock into school"),
        ("4", "Co-issue / issue", "Method + reference"),
        ("5", "Still owed", "Finish kits"),
        ("6", "Reports + print", "Both portals"),
    ]
    for i, (n, title, sub) in enumerate(steps):
        x = 0.4 + i * 2.15
        card = add_rect(s, Inches(x), Inches(2.3), Inches(2.0), Inches(3.6), fill=TEAL_SOFT)
        add_rect(s, Inches(x + 0.65), Inches(2.6), Inches(0.7), Inches(0.7), fill=TEAL)
        add_text_box(
            s, Inches(x + 0.65), Inches(2.72), Inches(0.7), Inches(0.5),
            n, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        )
        add_text_box(
            s, Inches(x + 0.15), Inches(3.55), Inches(1.7), Inches(1.2),
            title, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
        )
        add_text_box(
            s, Inches(x + 0.15), Inches(4.9), Inches(1.7), Inches(0.7),
            sub, size=11, color=MUTED, align=PP_ALIGN.CENTER,
        )
        add_appear_animation(s, card, i + 1)

    track(
        screenshot_slide(
            prs, "04  GETTING STARTED", "Sign in", "um-01-login",
            "Glass sign-in over Kenyan senior school uniforms",
            ["Open /login", "Use role email + desk1234", "Theme menu top-left", "Demo accounts under the form"],
        )
    )

    track(
        divider_slide(
            prs, "CHAPTER A", "Supplier portal",
            [
                "Own the catalog, schools, deliveries, and invoices",
                "Co-issue on admission day with school staff",
                "View issued uniforms and school stock (read-only)",
            ],
        ),
        footer=False,
    )

    supplier_shots = [
        ("A1  HOME", "Supplier home", "um-30-supplier-home", "Overview & quick actions",
         ["Co-issue desk CTA", "Reports & schools shortcuts", "Linked-school portfolio stats"]),
        ("A2  SCHOOLS", "Create & link schools", "um-31-supplier-schools", "Supplier admin only",
         ["Create school = name + code", "Adds school reporter login", "Or link existing code (GFS/RVA)", "Open Co-issue from cards"]),
        ("A3  CATALOG", "Supplier catalog", "um-32-supplier-catalog", "SKUs must match school items",
         ["Add products & sizes", "Unit price for invoices", "SKU is the receive key"]),
        ("A4  CO-ISSUE", "Co-issue desk", "um-33-supplier-coissue", "Issue against linked school stock",
         ["Pick school if more than one", "New student or find roster", "Load kit / lines", "Payment method + reference"]),
        ("A5  STILL OWED", "Incomplete kits", "um-34-supplier-still-owed", "Follow up short issues",
         ["List per linked school", "Issue what’s left", "Print list"]),
        ("A6  REPORTS", "Issued today", "um-35-supplier-reports", "Daily issue visibility",
         ["Issued list with payment refs", "Print report", "Jump to Still owed"]),
        ("A7  REPORTS", "Stock on hand (read-only)", "um-36-supplier-reports-stock", "School does stock take",
         ["Campus balances by SKU/size", "Low stock highlighted", "Print current view"]),
        ("A8  ORDERS", "Supply orders", "um-37-supplier-orders", "Create POs for linked schools",
         ["Order on behalf of school", "Confirm → pack into DN"]),
        ("A9  DELIVERIES", "Delivery notes", "um-38-supplier-deliveries", "Pack & dispatch",
         ["Create/pack DN", "Mark in transit", "Print DN on detail"]),
        ("A10  DN DETAIL", "Printable delivery note", "um-44-supplier-delivery-detail", "Operational DN pack",
         ["Print DN button", "Create invoice when ready", "School receives into stock"]),
        ("A11  INVOICES", "Invoices", "um-39-supplier-invoices", "Bill & collect",
         ["Invoice from delivery", "Mark paid / M-Pesa sandbox", "Print invoice"]),
        ("A12  INVOICE DETAIL", "Printable invoice", "um-45-supplier-invoice-detail", "Settlement record",
         ["Lines & totals", "Payment history", "Print invoice"]),
        ("A13  ACTIVITY", "Activity feed", "um-40-supplier-activity", "Orders · DN · invoices · co-issue",
         ["Timeline for support", "Correlation IDs (PO/DN/INV)"]),
        ("A14  BRANDING", "White-label branding", "um-43-supplier-branding", "Supplier admin",
         ["Brand name, mark, color", "Support contacts", "Portal accent updates"]),
        ("A15  STAFF", "Supplier staff home", "um-50-supplier-staff-home", "Ops without branding",
         ["Same co-issue & deliveries", "No branding menu"]),
    ]
    for row in supplier_shots:
        track(screenshot_slide(prs, *row))

    track(
        divider_slide(
            prs, "CHAPTER B", "School reporter desk",
            [
                "Co-issue at admission with supplier staff",
                "Stock take / adjust and receive delivery notes",
                "Reports, still owed, and printable lists",
            ],
        ),
        footer=False,
    )

    school_shots = [
        ("B1  HOME", "School desk home", "um-10-school-home", "Today at a glance",
         ["Issued / shortages / low stock", "Jump to issue, stock, reports", "Needs-attention notices"]),
        ("B2  ISSUE", "Issue desk", "um-11-school-issue", "Admission uniform issue",
         ["New student or find existing", "Load kit (requirements)", "Payment method + reference", "No parent signature"]),
        ("B3  STILL OWED", "Still to receive", "um-12-school-still-owed", "Incomplete uniforms",
         ["Auto-created on shortage", "Issue what’s left", "Print list"]),
        ("B4  STOCK", "Stock balances", "um-13-school-stock", "Stock take & adjust",
         ["Adjust with reason (Stock take)", "Ledger of receive/issue/void", "Print stock"]),
        ("B5  REPORTS", "School reports", "um-14-school-reports", "Issued today & shortages",
         ["Payment method shown", "Audit CSV export", "Print report"]),
        ("B6  STUDENTS", "Student roster", "um-15-school-students", "Admission master",
         ["Add / import CSV", "Open student history"]),
        ("B7  HISTORY", "Student history", "um-21-school-student-history", "Issue trail per student",
         ["Slips & shortages", "Useful for desk queries"]),
        ("B8  SLIP", "Issue record", "um-22-school-slip", "Internal audit slip",
         ["Issuer & payment ref", "Void restores stock + owed"]),
        ("B9  DELIVERIES", "Inbound deliveries", "um-16-school-deliveries", "Receive supplier DN",
         ["Open DN → Receive into stock", "SKU must match school items"]),
        ("B10  DN DETAIL", "Receive DN", "um-23-school-delivery-detail", "Posts ledger + balances",
         ["Optional receive note", "Print DN available"]),
        ("B11  RECEIVE", "Manual receive", "um-17-school-receive", "Non-DN inbound",
         ["Use when stock arrives without DN"]),
        ("B12  ACTIVITY", "School activity", "um-18-school-activity", "Issue & stock timeline",
         ["Desk support & audit"]),
        ("B13  SEARCH", "Desk search", "um-20-school-search", "Find students & slips",
         ["Top bar or Ops → Search"]),
    ]
    for row in school_shots:
        track(screenshot_slide(prs, *row))

    # Print rhythm
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    section_banner(s, "05  PRINT & DAILY RHYTHM", "What to print · when")
    tips = [
        ("Morning", "School: Still owed & low stock. Supplier: Reports → Issued."),
        ("Admission", "Co-issue or school Issue — method + reference; note shortages."),
        ("Receiving", "School Deliveries → Receive DN → stock rises."),
        ("Stock take", "School Stock → Adjust “Stock take” → Print stock."),
        ("Close of day", "Print Reports. Supplier may print stock view."),
        ("Documents", "DN & Invoice detail → Print DN / Print invoice."),
    ]
    for i, (t, body) in enumerate(tips):
        y = 1.65 + i * 0.85
        card = add_rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(0.75), fill=TEAL_SOFT if i % 2 == 0 else CREAM)
        add_text_box(s, Inches(0.75), Inches(y + 0.12), Inches(2.2), Inches(0.5), t, size=14, bold=True, color=TEAL)
        add_text_box(s, Inches(3.1), Inches(y + 0.18), Inches(9.3), Inches(0.5), body, size=13, color=NAVY)
        add_appear_animation(s, card, i + 1)

    # Appendix
    s = track(prs.slides.add_slide(prs.slide_layouts[6]))
    enable_transition(s)
    section_banner(s, "06  APPENDIX", "Demo credentials & routes")
    add_bullets(
        s, Inches(0.7), Inches(1.75), Inches(5.8), Inches(5.0),
        [
            "Password: desk1234 (all seed users)",
            "Supplier admin → /supplier",
            "Create school → /supplier/schools",
            "Co-issue → /supplier/issue",
            "Supplier reports → /supplier/reports",
            "School desk → /",
            "Issue → /issue · Still owed → /incomplete",
            "Stock → /stock · Reports → /reports",
        ],
        size=14,
    )
    add_bullets(
        s, Inches(7.0), Inches(1.75), Inches(5.5), Inches(5.0),
        [
            "School codes (seed): GFS, RVA",
            "Kits define admission requirements",
            "SKU match required for DN receive",
            "Payment = method + reference only",
            "Retired school purchase UI → home",
            "Rebuild: capture-user-manual.mjs",
            "Then: build-user-manual-pptx.py",
            "Also see docs/DESK_GUIDE.md",
        ],
        size=14,
    )

    # Close
    s = track(prs.slides.add_slide(prs.slide_layouts[6]), footer=False)
    enable_transition(s)
    add_bg(s, TEAL)
    add_text_box(
        s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(0.8),
        "You’re ready to run UniformDesk", size=32, bold=True, color=WHITE,
    )
    add_text_box(
        s, Inches(0.8), Inches(3.5), Inches(11), Inches(1.2),
        "Create schools · Co-issue at admission · Finish still owed · Print reports.\n"
        "Use Activity feeds and issue slips for the audit trail.",
        size=16, color=TEAL_SOFT,
    )

    total = len(slides)
    for i, (slide, want_footer) in enumerate(slides, start=1):
        if want_footer:
            add_footer(slide, i, total)

    prs.save(OUT)
    print(f"Wrote {OUT} ({total} slides)")
    missing = [p.stem for p in [] ]
    # report missing shots referenced
    needed = set()
    for p in SHOTS.glob("um-*.png"):
        needed.add(p.name)
    print(f"Screenshot files present: {len(list(SHOTS.glob('um-*.png')))}")
    return total


if __name__ == "__main__":
    build()
